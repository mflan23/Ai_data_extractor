"""
Document Processing Pipeline
Handles text extraction from various document formats including PDF, DOCX, and images with OCR.
"""

import os
import tempfile
from typing import Dict, Any, Optional
from pathlib import Path


def extract_text_from_pdf(filepath: str) -> str:
    """Extract text from PDF files using pdfplumber"""
    try:
        import pdfplumber
        text = ""
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    except ImportError:
        raise ImportError("pdfplumber is required. Install with: uv add pdfplumber")
    except Exception as e:
        raise Exception(f"Failed to extract text from PDF: {str(e)}")


def extract_text_from_docx(filepath: str) -> str:
    """Extract text from DOCX files using python-docx"""
    try:
        from docx import Document
        doc = Document(filepath)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except ImportError:
        raise ImportError("python-docx is required. Install with: uv add python-docx")
    except Exception as e:
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")


def extract_text_from_image(filepath: str) -> str:
    """Extract text from images using OCR (placeholder for Tesseract)"""
    # This would integrate with Tesseract or pytesseract
    # For now, return a placeholder message
    return f"[OCR Placeholder] Image file detected: {os.path.basename(filepath)}\n\n" \
           f"Note: Tesseract OCR integration required for actual text extraction.\n" \
           f"Install with: uv add pytesseract && pip install pytesseract"


def extract_text_from_txt(filepath: str) -> str:
    """Extract text from plain text files"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read().strip()
    except Exception as e:
        raise Exception(f"Failed to read text file: {str(e)}")


def extract_text_from_pptx(filepath: str) -> str:
    """Extract text from PowerPoint files"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text + "\n"
        return text.strip()
    except ImportError:
        raise ImportError("python-pptx is required. Install with: uv add python-pptx")
    except Exception as e:
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")


def detect_file_type(filepath: str) -> str:
    """Detect file type based on extension"""
    ext = Path(filepath).suffix.lower()
    type_map = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'docx',
        '.txt': 'txt',
        '.jpg': 'image',
        '.jpeg': 'image',
        '.png': 'image',
        '.gif': 'image',
        '.bmp': 'image',
        '.tiff': 'image',
        '.pptx': 'pptx',
        '.ppt': 'pptx',
    }
    return type_map.get(ext, 'unknown')


def process_document(filepath: str) -> Dict[str, Any]:
    """
    Main document processing function.
    Detects file type and extracts text accordingly.
    
    Args:
        filepath: Path to the document file
        
    Returns:
        Dictionary with extraction results
    """
    file_type = detect_file_type(filepath)
    filename = os.path.basename(filepath)
    
    print(f"Processing {filename}...")
    print(f"Detected type: {file_type}")
    
    try:
        if file_type == 'pdf':
            text = extract_text_from_pdf(filepath)
        elif file_type == 'docx':
            text = extract_text_from_docx(filepath)
        elif file_type == 'image':
            text = extract_text_from_image(filepath)
        elif file_type == 'txt':
            text = extract_text_from_txt(filepath)
        elif file_type == 'pptx':
            text = extract_text_from_pptx(filepath)
        else:
            return {
                'success': False,
                'filename': filename,
                'error': f'Unsupported file type: {file_type}',
                'raw_text': ''
            }
        
        return {
            'success': True,
            'filename': filename,
            'file_type': file_type,
            'raw_text': text,
            'word_count': len(text.split()),
            'char_count': len(text)
        }
        
    except Exception as e:
        return {
            'success': False,
            'filename': filename,
            'error': str(e),
            'raw_text': ''
        }


def cleanup_temp_file(filepath: str) -> None:
    """Clean up temporary files"""
    try:
        if os.path.exists(filepath):
            os.remove(filepath)
    except Exception as e:
        print(f"Warning: Failed to cleanup temp file: {str(e)}")


if __name__ == "__main__":
    # Test the processor
    import sys
    
    if len(sys.argv) > 1:
        filepath = sys.argv[1]
        result = process_document(filepath)
        print(f"\nResult: {result}")
    else:
        print("Usage: python document_processor.py <filepath>")
